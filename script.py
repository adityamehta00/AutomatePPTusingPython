import openpyxl
from pptx import Presentation
import json

#configuring from json file#
with open('configuration.json') as json_data:
    config = json.load(json_data)
    json_data.close()

#assigning value from directory (json file)
Pres_Template = config['TemplatePPT']
LayoutTestOutputExcel = config['LayoutTestOutputExcel']
No_of_placeholders_onlyName = config['No_of_placeholders_onlyName']
No_of_placeholders_double = config['No_of_placeholders_double']
No_of_placeholders_single =config['No_of_placeholders_single']


#opening the excel with all data
wb = openpyxl.load_workbook('data.xlsx')
ws = wb.active

#opening excel with placeholder ids
w = openpyxl.load_workbook(LayoutTestOutputExcel)
onlyName = w.get_sheet_by_name('onlyName')
double = w.get_sheet_by_name('double')
single = w.get_sheet_by_name('single')

#ppt with the master templates
prs = Presentation(Pres_Template)
layout_onlyName = 0         #the order is as they appear in powerpoint file
layout_double = 1
layout_single = 2
slide_layout_onlyName = prs.slide_layouts[layout_onlyName]
slide_layout_double = prs.slide_layouts[layout_double]
slide_layout_single = prs.slide_layouts[layout_single]


#giving index to row 
Num_header = 1
row = Num_header + 1


## FOR ONLY NAMES ##

print("Begining adding only names")
counter = 0
i = 0
#Creating Local array to address to placeholder All values will be defined autmatically using Layout test output excel#
k =[ ] 
z = 1
while(z<=No_of_placeholders_onlyName):
    k.append(onlyName.cell(row = z, column = 2).value)
    z +=1
print("k assigned placeholder ids")


while ws.cell(row=row, column=2).value == 3:            #checking layout
    if ws.cell(row=row, column=3).value == "p":           #checking if student is present 
        name = ws.cell(row=row, column = 1).value
   
        if counter % 7 == 0:         ##adding new slide for every 8th student                                             
            slide = prs.slides.add_slide(slide_layout_onlyName)         #addig new slide
            slide.placeholders[k[0]].text = str(name)                             #adding name to slide in the first text box
            i = 1
            
        else:                   ##adding names on same slide
            slide.placeholders[k[i]].text = str(name)                   
            i +=1

    row +=1
    counter +=1

print ("Only names printed, with %d names printed" %(counter))
 

## FOR TWO STUDENTS PER SLIDE ##

print("Begining adding two students per slide")
counter = 1
#Creating Local array to address to placeholder for LEFT SIDE.
k1 =[ ] 
z = 1
while(z<=(No_of_placeholders_double/2)):          ### Assuming that both left and right side have same no. of placeholders
    k1.append(double.cell(row = z, column = 2).value)
    z +=1
print("k1 assigned placeholder ids")

#Creating Local array to address to placeholder for RIGHT SIDE. 
k2 =[ ] 
while(z <=No_of_placeholders_double):
    k2.append(double.cell(row = z, column = 2).value)
    z +=1
print("k2 assigned placeholder ids")


while ws.cell(row=row, column=2).value == 2:          #checking layout
    if ws.cell(row=row, column=3).value == "p":         #checking if student is present
        
        #assigning all values from the data excel to local variables.#
        name = ws.cell(row=row, column=1).value
        sb = ws.cell(row=row, column=4).value
        sb_pcm = ws.cell(row=row, column=5).value
        main = ws.cell(row=row, column=6).value
        main_cat = ws.cell(row=row, column=7).value
        adv = ws.cell(row=row, column=8).value
        adv_cat = ws.cell(row=row, column=9).value
        bits = ws.cell(row=row, column=10).value
        photo = ws.cell(row=row, column=11).value

        # for adding new slide for every odd student
        if counter % 2 != 0:           
            slide = prs.slides.add_slide(slide_layout_double) #slide added
            
            slide.placeholders[k1[0]].insert_picture(photo)     #photo added
            slide.placeholders[k1[1]].text = str(name)             #name added

            i=2
            if sb !=0:  #this is used to check if we have the required data or not. If value = 0 then the correcpoding field will not be added to the slide.
                slide.placeholders[k1[i]].text = 'SB % :'
                slide.placeholders[k1[i+1]].text = str(sb)
                i += 2    
            if sb_pcm !=0:
                slide.placeholders[k1[i]].text = 'SB PCM% :'
                slide.placeholders[k1[i+1]].text = str(sb_pcm)
                i +=2     
            if main !=0:
                slide.placeholders[k1[i]].text = 'Mains Rank :'
                slide.placeholders[k1[i+1]].text = 'AIR ' + str(main)
                i += 2
            if main_cat !=0:
                slide.placeholders[k1[i]].text = 'Mains Cat Rank :'
                slide.placeholders[k1[i+1]].text ='AIR ' + str(main_cat)
                i += 2
            if adv !=0:
                slide.placeholders[k1[i]].text = 'JEE Adv. Rank :'
                slide.placeholders[k1[i+1]].text ='AIR ' + str(adv)
                i += 2
            if adv_cat !=0:
                slide.placeholders[k1[i]].text = 'JEE Adv. Cat Rank :'
                slide.placeholders[k1[i+1]].text = 'AIR ' + str(adv_cat)
                i += 2
            if bits !=0:
                slide.placeholders[k1[i]].text = 'BITSAT Score:'
                slide.placeholders[k1[i+1]].text = str(bits)
                i += 2
            counter +=1
        # for adding every even student on other side on an exsisting slide 
        else:    
            slide.placeholders[k2[0]].insert_picture(photo)
            slide.placeholders[k2[1]].text = name

            i=2
            if sb !=0:
                slide.placeholders[k2[i]].text = 'SB % :'
                slide.placeholders[k2[i+1]].text = str(sb)
                i += 2    
            if sb_pcm !=0:
                slide.placeholders[k2[i]].text = 'SB PCM% :'
                slide.placeholders[k2[i+1]].text = str(sb_pcm)
                i +=2     
            if main !=0:
                slide.placeholders[k2[i]].text = 'Mains Rank :'
                slide.placeholders[k2[i+1]].text = 'AIR ' + str(main)
                i += 2
            if main_cat !=0:
                slide.placeholders[k2[i]].text = 'Mains Cat Rank :'
                slide.placeholders[k2[i+1]].text ='AIR ' + str(main_cat)
                i += 2
            if adv !=0:
                slide.placeholders[k2[i]].text = 'JEE Adv. Rank :'
                slide.placeholders[k2[i+1]].text ='AIR ' + str(adv)
                i += 2
            if adv_cat !=0:
                slide.placeholders[k2[i]].text = 'JEE Adv. Cat Rank :'
                slide.placeholders[k2[i+1]].text = 'AIR ' + str(adv_cat)
                i += 2
            if bits !=0:
                slide.placeholders[k2[i]].text = 'BITSAT Score:'
                slide.placeholders[k2[i+1]].text = str(bits)
                i += 2
            counter +=1
    #moving to next row and starting loop again    
    row +=1
    
    
print ("Two students per slide printed, with %d names printed" %(counter))
    

## FOR SINGLE STUDENT PER SLIDE ##
print("Begining adding single student per slide")
counter = 1
##Local array to address to placeholder
k3 =[ ]
z = 1
while(z <=No_of_placeholders_single):
    k3.append(single.cell(row = z, column = 2).value)
    z +=1
print("k3 assigned placeholder ids")


while ws.cell(row=row, column=2).value == 1:      #checking layout
    if ws.cell(row=row, column=3).value == "p":     #for checking if student is present
        
        
        #assigning all values# 
        name = ws.cell(row=row, column=1).value
        sb = ws.cell(row=row, column=4).value
        sb_pcm = ws.cell(row=row, column=5).value
        main = ws.cell(row=row, column=6).value
        main_cat = ws.cell(row=row, column=7).value
        adv = ws.cell(row=row, column=8).value
        adv_cat = ws.cell(row=row, column=9).value
        bits = ws.cell(row=row, column=10).value
        picture = ws.cell(row=row, column=11).value
    
        slide = prs.slides.add_slide(slide_layout_single) #slide added        
        slide.placeholders[k3[0]].insert_picture(photo)   #photo added
        slide.placeholders[k3[1]].text = str(name)             #name added   

        i=2
        if sb !=0:
            slide.placeholders[k3[i]].text = 'SB % :'
            slide.placeholders[k3[i+1]].text = str(sb)
            i += 2    
        if sb_pcm !=0:
            slide.placeholders[k3[i]].text = 'SB PCM% :'
            slide.placeholders[k3[i+1]].text = str(sb_pcm)
            i +=2     
        if main !=0:
            slide.placeholders[k3[i]].text = 'Mains Rank :'
            slide.placeholders[k3[i+1]].text = 'AIR ' + str(main)
            i += 2
        if main_cat !=0:
            slide.placeholders[k3[i]].text = 'Mains Cat Rank :'
            slide.placeholders[k3[i+1]].text ='AIR ' + str(main_cat)
            i += 2
        if adv !=0:
            slide.placeholders[k3[i]].text = 'JEE Adv. Rank :'
            slide.placeholders[k3[i+1]].text ='AIR ' + str(adv)
            i += 2
        if adv_cat !=0:
            slide.placeholders[k3[i]].text = 'JEE Adv Cat. Rank :'
            slide.placeholders[k3[i+1]].text = 'AIR ' + str(adv_cat)
            i += 2
        if bits !=0:
            slide.placeholders[k3[i]].text = 'BITSAT Score:'
            slide.placeholders[k3[i+1]].text = str(bits)
            i += 2
    
    counter +=1
    row +=1

print ("One student per slide printed, with %d names printed" %(counter))    

prs.save('output.pptx')

