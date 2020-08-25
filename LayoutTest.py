from pptx import Presentation
from xlwt import Workbook
import  openpyxl
import json

#configuring from json file#
with open('configuration.json') as json_data:
    config = json.load(json_data)
    json_data.close()

#assigning value from directory (json file)
Pres_Template = config['TemplatePPT']
LayoutTestExcel = config['LayoutTestExcel']
slideID_onlyName = config['slideID_onlyName']
slideID_double = config['slideID_double']
slideID_single = config['slideID_single']
No_of_placeholders_onlyName = config['No_of_placeholders_onlyName']
No_of_placeholders_double = config['No_of_placeholders_double']
No_of_placeholders_single =config['No_of_placeholders_single']


#opening the excel with names of placeholders
w = openpyxl.load_workbook(LayoutTestExcel)
onlyName = w.get_sheet_by_name('onlyName')
double = w.get_sheet_by_name('double')
single = w.get_sheet_by_name('single')

#opening a new workbook which will be saved as LayoutTestOutput.xls
wb = Workbook()
sheet1 = wb.add_sheet('onlyName')
sheet2 = wb.add_sheet('double')
sheet3 = wb.add_sheet('single')


prs = Presentation(Pres_Template)


#for only names
print (" ")
print("Printing and storing placeholder ids for only names template")
slide = prs.slides.get(slideID_onlyName)   #getting the slide contaitng only names from the template 
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    row = 1
    while(row<=No_of_placeholders_onlyName):
        k = onlyName.cell(row = row, column = 1).value
        if(k==shape.name):
            sheet1.write(row-1,0,shape.name)
            sheet1.write(row-1,1,shape.placeholder_format.idx)
            row=row+1  
        else:
            row=row+1
    
#for two student per slide
print (" ")
print("Printing  and storing placeholder ids for 2 names template")

slide = prs.slides.get(slideID_double)   #getting the slide contatining two names from the template
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
    row = 1
    while(row<=No_of_placeholders_double):
        k = double.cell(row = row, column = 1).value
        if(k==shape.name):
            sheet2.write(row-1,0,shape.name)
            sheet2.write(row-1,1,shape.placeholder_format.idx)
            row=row+1  
        else:
            row=row+1

#for one student per slide
print(" ")
print("Printing and storing placeholder ids for only 1 name template")
slide = prs.slides.get(slideID_single)   #getting the slide containing single name from the template.
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

    row = 1
    while(row<=No_of_placeholders_single):
        k = single.cell(row = row, column = 1).value
        if(k==shape.name):
            sheet3.write(row-1,0,shape.name)
            sheet3.write(row-1,1,shape.placeholder_format.idx)
            row=row+1  
        else:
            row=row+1

wb.save('LayoutTestOutput.xls')

